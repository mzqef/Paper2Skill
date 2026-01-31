"""Document loaders for multiple file formats."""

from pathlib import Path
from typing import Union
from abc import ABC, abstractmethod

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class DocumentLoader(ABC):
    """Abstract base class for document loaders."""

    @abstractmethod
    def load(self, file_path: Union[str, Path]) -> str:
        """Load and extract text from a document."""
        pass


class PDFLoader(DocumentLoader):
    """Loader for PDF documents."""

    def load(self, file_path: Union[str, Path]) -> str:
        """Load text from a PDF file."""
        if PdfReader is None:
            raise ImportError("pypdf is required for PDF support. Install with: pip install pypdf")
        
        text_content = []
        pdf = PdfReader(str(file_path))
        
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_content.append(text)
        
        return "\n\n".join(text_content)


class WordLoader(DocumentLoader):
    """Loader for Word documents (.docx)."""

    def load(self, file_path: Union[str, Path]) -> str:
        """Load text from a Word document."""
        if DocxDocument is None:
            raise ImportError(
                "python-docx is required for Word support. Install with: pip install python-docx"
            )
        
        doc = DocxDocument(str(file_path))
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        return "\n\n".join(text_content)


class PowerPointLoader(DocumentLoader):
    """Loader for PowerPoint presentations (.pptx)."""

    def load(self, file_path: Union[str, Path]) -> str:
        """Load text from a PowerPoint presentation."""
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PowerPoint support. "
                "Install with: pip install python-pptx"
            )
        
        prs = Presentation(str(file_path))
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            text_content.append("\n".join(slide_text))
        
        return "\n\n".join(text_content)


class MarkdownLoader(DocumentLoader):
    """Loader for Markdown documents."""

    def load(self, file_path: Union[str, Path]) -> str:
        """Load text from a Markdown file."""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()


class MultiFormatLoader:
    """Loader that automatically detects and loads multiple document formats."""

    LOADERS = {
        '.pdf': PDFLoader,
        '.docx': WordLoader,
        '.pptx': PowerPointLoader,
        '.md': MarkdownLoader,
        '.markdown': MarkdownLoader,
        '.txt': MarkdownLoader,
    }

    @classmethod
    def load(cls, file_path: Union[str, Path]) -> str:
        """
        Load a document of any supported format.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Extracted text content from the document
            
        Raises:
            ValueError: If file format is not supported
        """
        path = Path(file_path)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        suffix = path.suffix.lower()
        
        if suffix not in cls.LOADERS:
            raise ValueError(
                f"Unsupported file format: {suffix}. "
                f"Supported formats: {', '.join(cls.LOADERS.keys())}"
            )
        
        loader_class = cls.LOADERS[suffix]
        loader = loader_class()
        
        return loader.load(path)
